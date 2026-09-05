"""Build the IziMom 8-slide PowerPoint briefing."""

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml
from pptx.util import Emu, Inches, Pt
from lxml import etree

OUT = r"c:\Users\tuyiz\Desktop\KaziniKazi\izimom\IziMom-briefing.pptx"

CORAL = RGBColor(0xEB, 0x57, 0x6E)
CORAL_DARK = RGBColor(0xC9, 0x3D, 0x54)
BLUSH = RGBColor(0xFF, 0xF1, 0xF4)
CREAM = RGBColor(0xFA, 0xF6, 0xF2)
INK = RGBColor(0x24, 0x18, 0x1C)
MUTED = RGBColor(0x6B, 0x55, 0x5C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x2A, 0x10, 0x16)
DARK2 = RGBColor(0x3B, 0x15, 0x20)
LINE = RGBColor(0xEA, 0xD8, 0xDC)
PINK = RGBColor(0xFF, 0xB7, 0xC2)
SOFT = RGBColor(0xF3, 0xD4, 0xDA)

NSMAP = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}


def set_fill(shape, rgb):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb
    shape.line.fill.background()


def set_fill_line(shape, fill, line=None, width=Pt(1)):
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = width


def round_rect(slide, l, t, w, h, fill, line=None, radius=0.12):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    set_fill_line(shape, fill, line)
    shape.adjustments[0] = radius
    return shape


def oval(slide, l, t, w, h, fill):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, l, t, w, h)
    set_fill(shape, fill)
    return shape


def tb(slide, l, t, w, h, text, size=14, bold=False, color=INK, align=PP_ALIGN.LEFT, font="Calibri", italic=False, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    try:
        tf._txBody.bodyPr.set("anchor", {MSO_ANCHOR.TOP: "t", MSO_ANCHOR.MIDDLE: "ctr", MSO_ANCHOR.BOTTOM: "b"}[anchor])
    except Exception:
        pass
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font
    return box


def bullets(slide, l, t, w, h, items, size=13, color=INK, bullet_color=CORAL):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = 0
        p.space_after = Pt(6)
        run = p.add_run()
        run.text = "•  " + item
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.name = "Calibri"
    return box


def footer(slide, n):
    tb(slide, Inches(0.5), Inches(7.18), Inches(8), Inches(0.22), "IziMom  ·  Confidential briefing  ·  Rwanda 2026", 10, False, MUTED)
    tb(slide, Inches(11.4), Inches(7.18), Inches(1.4), Inches(0.22), f"{n}  /  8", 10, True, CORAL, PP_ALIGN.RIGHT)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.44), Inches(13.333), Inches(0.06))
    set_fill(bar, CORAL)


def content_bg(slide):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    set_fill(bg, CREAM)
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.08), Inches(7.5))
    set_fill(accent, CORAL)


def eyebrow(slide, text, y=0.28):
    tb(slide, Inches(0.5), Inches(y), Inches(12), Inches(0.28), text.upper(), 11, True, CORAL)


def heading(slide, text, y=0.5, h=0.7):
    tb(slide, Inches(0.5), Inches(y), Inches(12.3), Inches(h), text, 26, True, INK, font="Georgia")


def sub(slide, text, y=1.15):
    tb(slide, Inches(0.5), Inches(y), Inches(12.3), Inches(0.42), text, 14, False, MUTED)


def card_title(slide, l, t, w, text, color=INK):
    tb(slide, l + Inches(0.16), t + Inches(0.1), w - Inches(0.3), Inches(0.32), text, 14, True, color)


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

# ---------------------------------------------------------------------------
# Slide 1 — Title
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(blank)
bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
set_fill(bg, DARK)
glow = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8.2), Inches(-1.4), Inches(7), Inches(4.4))
set_fill(glow, RGBColor(0x6B, 0x22, 0x34))
glow2 = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-1.6), Inches(4.6), Inches(5.2), Inches(3.6))
set_fill(glow2, RGBColor(0x4A, 0x18, 0x26))

tb(s, Inches(0.7), Inches(0.55), Inches(12), Inches(0.3), "IZIHOME FAMILY  ·  RWANDA  ·  AUGUST 2026", 12, True, PINK)

mark = oval(s, Inches(0.7), Inches(1.05), Inches(0.48), Inches(0.48), CORAL)
tb(s, Inches(1.3), Inches(1.1), Inches(4), Inches(0.42), "IziMom", 22, True, WHITE)

tb(s, Inches(0.7), Inches(2.0), Inches(12), Inches(0.7), "Rwanda won the birth.", 40, True, WHITE, font="Georgia")
tb(s, Inches(0.7), Inches(2.7), Inches(12), Inches(0.7), "Mothers are left after it.", 40, True, PINK, font="Georgia", italic=True)

tb(
    s,
    Inches(0.7),
    Inches(3.6),
    Inches(10.5),
    Inches(0.85),
    "Licensed nurses and midwives at home for the first 6–12 weeks after delivery — recovery, newborn care, feeding, rest, and emotional support.",
    16,
    False,
    SOFT,
)

pills = [
    (0.7, "The problem, with numbers"),
    (3.7, "Care packages"),
    (6.15, "6 nurses  ·  2 cities"),
    (9.15, "How we register"),
]
for x, label in pills:
    r = round_rect(s, Inches(x), Inches(4.75), Inches(2.7), Inches(0.42), DARK2, WHITE, 0.5)
    tb(s, Inches(x), Inches(4.8), Inches(2.7), Inches(0.34), label, 12, True, WHITE, PP_ALIGN.CENTER)

tb(s, Inches(0.7), Inches(6.85), Inches(12), Inches(0.28), "We care, so you can heal.", 13, False, PINK, italic=True)

# ---------------------------------------------------------------------------
# Slide 2 — Aftercare gap
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(blank)
content_bg(s)
eyebrow(s, "The aftercare gap")
heading(s, "Hospital care ends at discharge. Motherhood does not.")
sub(s, "Rwanda is a global success on facility birth. The missing product is the weeks at home — when most complications, feeding problems, and exhaustion actually happen.")

stats = [
    ("~396k", "live births in Rwanda in 2023", False),
    ("93%", "deliver in a health facility", False),
    ("70%", "get a PNC check in the first 2 days", False),
    ("44%", "receive adequate PNC content", True),
]
for i, (num, label, warn) in enumerate(stats):
    x = 0.5 + i * 3.15
    fill, nc, lc = (DARK2, PINK, SOFT) if warn else (WHITE, CORAL, MUTED)
    round_rect(s, Inches(x), Inches(1.7), Inches(3.0), Inches(1.45), fill, None if warn else LINE)
    tb(s, Inches(x + 0.16), Inches(1.82), Inches(2.7), Inches(0.7), num, 32, True, nc, font="Georgia")
    tb(s, Inches(x + 0.16), Inches(2.5), Inches(2.7), Inches(0.5), label, 12, False, lc)

round_rect(s, Inches(0.5), Inches(3.35), Inches(6.15), Inches(3.55), WHITE, LINE)
card_title(s, Inches(0.5), Inches(3.42), Inches(6.15), "What “adequate aftercare” actually means")
tb(
    s,
    Inches(0.7),
    Inches(3.85),
    Inches(5.75),
    Inches(0.95),
    "WHO wants four postnatal contacts: within 24 hours, days 3–4, days 7–14, and around 6 weeks. Rwanda’s first-two-days check is improving. Continuity after that is thin.",
    13,
    False,
    MUTED,
)
bullets(
    s,
    Inches(0.7),
    Inches(4.85),
    Inches(5.75),
    Inches(1.8),
    [
        "Only 55% of mothers are counselled on newborn danger signs",
        "Only 19% of infants had a postnatal check within 2 months",
        "Working mothers are less likely to complete adequate PNC",
    ],
    13,
)

round_rect(s, Inches(6.85), Inches(3.35), Inches(5.95), Inches(3.55), DARK2, CORAL)
card_title(s, Inches(6.85), Inches(3.42), Inches(5.95), "IziMom’s opening", WHITE)
tb(
    s,
    Inches(7.05),
    Inches(3.9),
    Inches(5.55),
    Inches(1.35),
    "The public system delivers the baby. Families still need a trained person in the house for bathing, feeding, recovery, meals, and spotting danger signs — especially first-time and employed mothers in cities.",
    14,
    False,
    WHITE,
)
tb(
    s,
    Inches(7.05),
    Inches(5.4),
    Inches(5.55),
    Inches(1.3),
    "Sources: NISR / Our World in Data births 2023; RDHS 2019–20 (93% facility, 94% skilled, 70% PNC in 2 days); Kawuki et al. 2022 (44.3% adequate PNC; 55.4% danger-sign counselling); Scientific Reports 2026 DHS analysis (18.7% infant PNC in 2 months).",
    10,
    False,
    SOFT,
)
footer(s, 2)

# ---------------------------------------------------------------------------
# Slide 3 — Depression & knowledge
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(blank)
content_bg(s)
eyebrow(s, "Depression and knowledge")
heading(s, "One in five mothers is struggling. Almost nobody is treated.")

nums = [
    ("1 in 5", "postnatal women in Rwanda report depressive symptoms (~21%). One in four is depressed during pregnancy."),
    ("5.3%", "of people with mental health needs use services. There is no national perinatal mental-health guideline."),
    ("88%", "of surveyed mothers had no internet — so an app alone cannot close the knowledge gap."),
]
for i, (num, label) in enumerate(nums):
    x = 0.5 + i * 4.2
    round_rect(s, Inches(x), Inches(1.35), Inches(4.0), Inches(1.85), WHITE, LINE)
    tb(s, Inches(x + 0.18), Inches(1.45), Inches(3.65), Inches(0.6), num, 32, True, CORAL, font="Georgia")
    tb(s, Inches(x + 0.18), Inches(2.1), Inches(3.65), Inches(0.95), label, 12, False, MUTED)

round_rect(s, Inches(0.5), Inches(3.4), Inches(6.15), Inches(3.5), WHITE, LINE)
card_title(s, Inches(0.5), Inches(3.48), Inches(6.15), "Mothers do not know what “not okay” looks like")
bullets(
    s,
    Inches(0.7),
    Inches(3.95),
    Inches(5.75),
    Inches(2.7),
    [
        "Low literacy on perinatal mental-health signs; women wait for a “fever” before seeking help",
        "Stigma: distress is treated as weakness, so families minimise it",
        "Exclusive breastfeeding fell from 87% (2015) to 81% (2020)",
        "Only 55% get counselling on newborn danger signs after birth",
    ],
    13,
)

round_rect(s, Inches(6.85), Inches(3.4), Inches(5.95), Inches(3.5), WHITE, LINE)
card_title(s, Inches(6.85), Inches(3.48), Inches(5.95), "What a nurse at home changes")
bullets(
    s,
    Inches(7.05),
    Inches(3.95),
    Inches(5.55),
    Inches(2.0),
    [
        "Teaches feeding, bathing, safe sleep, and warning signs in Kinyarwanda",
        "Screens mood on each visit and refers when needed",
        "Gives the mother rest — a proven buffer against postnatal depression",
        "Includes the partner, which studies show is a top protective factor",
    ],
    13,
)
tb(
    s,
    Inches(7.05),
    Inches(6.05),
    Inches(5.55),
    Inches(0.7),
    "Sources: Umuziga et al. 2021 / Front. Glob. Womens Health 2023 (20.9%); PMC 2025; Rwanda Mental Health Survey (5.3%); RDHS 2019–20.",
    10,
    False,
    MUTED,
)
footer(s, 3)

# ---------------------------------------------------------------------------
# Slide 4 — Solution
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(blank)
content_bg(s)
eyebrow(s, "The offer")
heading(s, "IziMom: a licensed nurse in the home for the fourth trimester.")
sub(s, "Same model as MomHelper — built for Rwanda, under the IziHome brand. We do not replace the health centre. We cover the hours and the skills the ward cannot.")

services = [
    ("Newborn care", "Bathing, soothing, safe sleep, weight and feeding cues."),
    ("Mother recovery", "Wound / C-section checks, rest coaching, when to return to the facility."),
    ("Feeding support", "Breastfeeding, latch, expressing, and when formula is needed."),
    ("Mood check-ins", "Listen, screen, destigmatise, and refer — not a psychiatry clinic."),
    ("Meals & light home", "Nourishing food for mum, baby laundry, a calmer room."),
    ("Night support", "Overnight watch so she can sleep — the intervention families ask for first."),
]
for i, (title, body) in enumerate(services):
    col, row = i % 3, i // 3
    x, y = 0.5 + col * 4.2, 1.7 + row * 1.85
    round_rect(s, Inches(x), Inches(y), Inches(4.0), Inches(1.7), WHITE, LINE)
    oval(s, Inches(x + 0.18), Inches(y + 0.2), Inches(0.28), Inches(0.28), CORAL)
    tb(s, Inches(x + 0.55), Inches(y + 0.18), Inches(3.25), Inches(0.35), title, 16, True, INK)
    tb(s, Inches(x + 0.18), Inches(y + 0.65), Inches(3.65), Inches(0.85), body, 13, False, MUTED)

round_rect(s, Inches(0.5), Inches(5.55), Inches(12.3), Inches(1.35), DARK2)
tb(
    s,
    Inches(0.7),
    Inches(5.75),
    Inches(11.9),
    Inches(0.95),
    "Who we serve first: first-time mothers, C-section recoveries, twins, employed urban households, and families whose relatives live far away.",
    16,
    True,
    WHITE,
)
footer(s, 4)

# ---------------------------------------------------------------------------
# Slide 5 — Packages
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(blank)
content_bg(s)
eyebrow(s, "Packages")
heading(s, "Start cheap. Upgrade when the house needs more hands.")
sub(s, "Priced for Kigali and Huye households. All visits by NCNM-licensed nurses or midwives. WhatsApp line included on every plan.")

pkgs = [
    (
        "Weekly Essential",
        "10,000",
        "RWF / week",
        "First weeks at home",
        ["4 visits · 3 hours each", "Newborn bath & feeding", "Light cleaning of baby spaces", "WhatsApp check-ins"],
        False,
    ),
    (
        "Comfort Monthly",
        "38,000",
        "RWF / month",
        "Weekday recovery support",
        ["5 daytime visits each week", "Meals for mum", "Newborn care & laundry", "Nurse guidance on recovery"],
        True,
    ),
    (
        "Premium Recovery",
        "75,000",
        "RWF / month",
        "Day + nights",
        ["Daytime care + 2 overnights", "Dedicated midwife check-ins", "Breastfeeding & recovery coaching", "Priority on-demand visits"],
        False,
    ),
]
for i, (name, price, unit, who, items, feat) in enumerate(pkgs):
    x = 0.5 + i * 4.2
    fill, tc, pc, mc = (DARK2, WHITE, PINK, SOFT) if feat else (WHITE, INK, CORAL, MUTED)
    round_rect(s, Inches(x), Inches(1.65), Inches(4.0), Inches(4.55), fill, CORAL if feat else LINE, 0.08)
    if feat:
        tag = round_rect(s, Inches(x + 2.15), Inches(1.78), Inches(1.65), Inches(0.28), CORAL, None, 0.5)
        tb(s, Inches(x + 2.15), Inches(1.78), Inches(1.65), Inches(0.28), "Most chosen", 10, True, WHITE, PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    tb(s, Inches(x + 0.22), Inches(1.82), Inches(3.5), Inches(0.35), name, 16, True, tc)
    tb(s, Inches(x + 0.22), Inches(2.2), Inches(3.55), Inches(0.5), price, 28, True, pc, font="Georgia")
    tb(s, Inches(x + 0.22), Inches(2.7), Inches(3.55), Inches(0.28), unit, 12, False, mc)
    tb(s, Inches(x + 0.22), Inches(3.02), Inches(3.55), Inches(0.3), who, 12, False, mc)
    bullets(s, Inches(x + 0.22), Inches(3.4), Inches(3.55), Inches(2.4), items, 13, WHITE if feat else INK)

round_rect(s, Inches(0.5), Inches(6.32), Inches(12.3), Inches(0.72), WHITE, CORAL)
tb(s, Inches(0.7), Inches(6.4), Inches(8.5), Inches(0.28), "Night Support add-on", 14, True, INK)
tb(s, Inches(0.7), Inches(6.68), Inches(8.5), Inches(0.28), "Add overnight newborn care to any package.", 12, False, MUTED)
tb(s, Inches(9.3), Inches(6.48), Inches(3.3), Inches(0.42), "18,000 RWF / week", 16, True, CORAL, PP_ALIGN.RIGHT)
footer(s, 5)

# ---------------------------------------------------------------------------
# Slide 6 — Team & locations
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(blank)
content_bg(s)
eyebrow(s, "Launch team")
heading(s, "Six licensed nurses. Two cities. One protocol.")
sub(s, "Hire for coverage, not a clinic. Each nurse can carry about 8–12 Comfort families or a mix of weekly visits. We start where births and ability to pay concentrate.")

round_rect(s, Inches(0.5), Inches(1.7), Inches(6.15), Inches(3.55), WHITE, LINE)
tb(s, Inches(0.7), Inches(1.85), Inches(5.75), Inches(0.28), "KIGALI  ·  GASABO & KICUKIRO", 11, True, CORAL)
tb(s, Inches(0.7), Inches(2.2), Inches(5.75), Inches(0.55), "4 nurses", 32, True, INK, font="Georgia")
bullets(
    s,
    Inches(0.7),
    Inches(2.85),
    Inches(5.75),
    Inches(1.3),
    [
        "1 senior midwife (team lead)",
        "2 registered nurses — days",
        "1 pediatric / night nurse",
    ],
    14,
)
tb(
    s,
    Inches(0.7),
    Inches(4.25),
    Inches(5.75),
    Inches(0.8),
    "Catchment: CHUK, King Faisal, private maternities, Kimironko–Kicukiro corridor. Highest density of paying first-time mothers.",
    12,
    False,
    MUTED,
)

round_rect(s, Inches(6.85), Inches(1.7), Inches(5.95), Inches(3.55), WHITE, LINE)
tb(s, Inches(7.05), Inches(1.85), Inches(5.55), Inches(0.28), "HUYE  ·  SOUTHERN PROVINCE", 11, True, CORAL)
tb(s, Inches(7.05), Inches(2.2), Inches(5.55), Inches(0.55), "2 nurses", 32, True, INK, font="Georgia")
bullets(
    s,
    Inches(7.05),
    Inches(2.85),
    Inches(5.55),
    Inches(1.0),
    [
        "1 certified midwife",
        "1 maternal-health nurse",
    ],
    14,
)
tb(
    s,
    Inches(7.05),
    Inches(4.15),
    Inches(5.55),
    Inches(0.9),
    "Catchment: CHUB and university families. Proves the model outside Kigali without splitting the team across the whole country.",
    12,
    False,
    MUTED,
)

hire = [
    ("Must-haves", "NCNM licence · A1/A0 nursing or midwifery · clean criminal record · Kinyarwanda + English or French"),
    ("We train", "Home etiquette, EPDS mood screen, breastfeeding refresh, IziMom uniform, WhatsApp reporting"),
    ("Pay", "Competitive monthly + visit bonus · RSSB · phone & transport stipend"),
]
for i, (title, body) in enumerate(hire):
    x = 0.5 + i * 4.2
    round_rect(s, Inches(x), Inches(5.4), Inches(4.0), Inches(1.5), WHITE, LINE)
    tb(s, Inches(x + 0.16), Inches(5.5), Inches(3.7), Inches(0.3), title, 13, True, CORAL)
    tb(s, Inches(x + 0.16), Inches(5.82), Inches(3.7), Inches(0.95), body, 12, False, INK)
footer(s, 6)

# ---------------------------------------------------------------------------
# Slide 7 — Marketing
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(blank)
content_bg(s)
eyebrow(s, "Go to market")
heading(s, "Mothers do not search “postpartum nurse.” They ask the ward and WhatsApp.", h=0.85)

mkt = [
    ("01", "Hospital gate", "MOUs with maternity wards: CHUK, King Faisal, CHUB, polyclinics. Discharge flyer + nurse intro 24 hours before going home. Gift a first Essential week to C-section mothers as a trial."),
    ("02", "WhatsApp & radio", "Book in Kinyarwanda on WhatsApp. Short radio spots on Flash FM / Isango on danger signs and rest. 88% of mothers are not on the internet — radio still wins."),
    ("03", "Trusted women", "Pay CHWs and independent midwives a referral fee. Umugoroba w’ababyeyi talks. Church and Sacco women’s groups. Instagram/TikTok only as social proof, not the main funnel."),
    ("04", "Employers", "Sell Comfort as a 6-week return-to-work benefit to banks, NGOs, and hotels. One HR deal fills a nurse’s roster. Cross-sell from IziFresh, IziWash, IziFix."),
]
for i, (num, title, body) in enumerate(mkt):
    col, row = i % 2, i // 2
    x, y = 0.5 + col * 6.35, 1.55 + row * 2.15
    round_rect(s, Inches(x), Inches(y), Inches(6.15), Inches(2.0), WHITE, LINE)
    oval(s, Inches(x + 0.18), Inches(y + 0.2), Inches(0.42), Inches(0.42), BLUSH)
    tb(s, Inches(x + 0.18), Inches(y + 0.24), Inches(0.42), Inches(0.36), num, 11, True, CORAL, PP_ALIGN.CENTER)
    tb(s, Inches(x + 0.72), Inches(y + 0.22), Inches(5.2), Inches(0.38), title, 16, True, INK)
    tb(s, Inches(x + 0.22), Inches(y + 0.72), Inches(5.7), Inches(1.15), body, 13, False, MUTED)

round_rect(s, Inches(0.5), Inches(5.95), Inches(12.3), Inches(0.95), DARK2)
tb(
    s,
    Inches(0.7),
    Inches(6.1),
    Inches(11.9),
    Inches(0.7),
    "First 90 days KPI: 40 paying families, 60% from hospital/midwife referrals, NPS asked at week 2. Spend on people at the ward door, not ads.",
    15,
    True,
    WHITE,
)
footer(s, 7)

# ---------------------------------------------------------------------------
# Slide 8 — Registration
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(blank)
content_bg(s)
eyebrow(s, "Legal path")
heading(s, "Register IziMom Ltd, then licence the care — not the other way around.", h=0.8)

regs = [
    ("1", "RDB company", "Online at org.rdb.rw. Free. ~6 hours. Need: approved name, national ID/passport, email, physical address, activities (home nursing / postpartum care), shareholders and directors. Certificate of incorporation."),
    ("2", "RRA tax", "TIN comes with RDB. Register with RRA within 7 days of starting. VAT later, when turnover exceeds 20 million RWF in 12 months (or 5 million in 3 consecutive months)."),
    ("3", "MoH facility licence", "Irembo: Provisional authorisation (free, ~30 days) then operating licence after inspection. Attach: RDB certificate, TIN, MD CV, 5-year plan, MD criminal-record certificate, site UPI and capacity."),
    ("4", "Nurses + RSSB", "Every nurse/midwife must hold a current National Council of Nurses and Midwives practising licence. Register as an employer with RSSB before the first payroll."),
]
for i, (num, title, body) in enumerate(regs):
    col, row = i % 2, i // 2
    x, y = 0.5 + col * 6.35, 1.45 + row * 1.85
    round_rect(s, Inches(x), Inches(y), Inches(6.15), Inches(1.7), WHITE, LINE)
    oval(s, Inches(x + 0.18), Inches(y + 0.18), Inches(0.38), Inches(0.38), BLUSH)
    tb(s, Inches(x + 0.18), Inches(y + 0.22), Inches(0.38), Inches(0.32), num, 13, True, CORAL, PP_ALIGN.CENTER)
    tb(s, Inches(x + 0.68), Inches(y + 0.2), Inches(5.2), Inches(0.35), title, 15, True, INK)
    tb(s, Inches(x + 0.22), Inches(y + 0.62), Inches(5.7), Inches(0.95), body, 12, False, MUTED)

round_rect(s, Inches(0.5), Inches(5.3), Inches(6.15), Inches(1.6), WHITE, LINE)
tb(s, Inches(0.7), Inches(5.42), Inches(5.75), Inches(0.3), "90-day sequence", 14, True, INK)
tb(
    s,
    Inches(0.7),
    Inches(5.78),
    Inches(5.75),
    Inches(0.95),
    "Week 1–2 RDB + TIN + bank. Week 3–6 MoH file + hire 6. Week 6–8 training and uniforms. Week 8 first paid visits in Kigali. Month 3 Huye live.",
    13,
    False,
    MUTED,
)

round_rect(s, Inches(6.85), Inches(5.3), Inches(5.95), Inches(1.6), DARK2, CORAL)
tb(s, Inches(7.05), Inches(5.42), Inches(5.55), Inches(0.3), "What we are asking", 14, True, WHITE)
tb(
    s,
    Inches(7.05),
    Inches(5.78),
    Inches(5.55),
    Inches(0.95),
    "Approval to incorporate IziMom Ltd, seed for 6 salaries and 3 months of operations, and introductions to CHUK and CHUB maternity leads.",
    13,
    False,
    WHITE,
)
footer(s, 8)

prs.save(OUT)
print(OUT)
